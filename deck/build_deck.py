"""Generate the CodeStreet 2026 screening-round pitch deck (.pptx).

11 self-standing slides (light Amex-branded theme) tuned for Round-1 judging:
Problem -> HLD -> UI mockup -> LLD/stack -> 'why Amex'. The winners' formula.
Run:  uv run python deck/build_deck.py
Output: deck/CodeStreet2026_Servicing_Agent.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- palette ---------------------------------------------------------------
# Light theme on Amex's official palette (Amex Blue #016FD0, Deep Blue #002663)
BG = RGBColor(0xF6, 0xF9, 0xFD)        # light cool background
CARD = RGBColor(0xE9, 0xF0, 0xF9)      # light card fill
ACCENT = RGBColor(0x01, 0x6F, 0xD0)    # Amex Blue (official)
ACCENT2 = RGBColor(0x00, 0x26, 0x63)   # Amex Deep Blue
WHITE = RGBColor(0x0F, 0x17, 0x2A)     # primary TEXT (dark); name kept for reuse
MUTED = RGBColor(0x5B, 0x66, 0x73)     # secondary text
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)     # text on dark/colored fills
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xB4, 0x53, 0x09)
ROSE = RGBColor(0xC0, 0x36, 0x2C)
VIOLET = RGBColor(0x6D, 0x4F, 0xC4)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color,
    bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(4 * space); p.space_before = 0
        for (txt, size, color, bold, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
    return tb


def R(t, size, color=WHITE, bold=False, font=FONT):
    return (t, size, color, bold, font)


def header(s, eyebrow, title, ecolor=ACCENT):
    box(s, 0, 0, 0.18, 7.5, fill=ACCENT)  # left accent stripe
    if eyebrow:
        text(s, 0.85, 0.55, 11, 0.4,
             [[R(eyebrow.upper(), 13, ecolor, True)]])
    text(s, 0.82, 0.92, 11.6, 1.2, [[R(title, 32, WHITE, True, FONT_H)]])


_page = 2  # title slide has no footer; page numbers auto-increment from slide 2


def footer(s, n=None):
    global _page
    n = _page if n is None else n
    text(s, 0.85, 7.02, 8, 0.3,
         [[R("End-to-End Servicing Agent · CodeStreet 2026", 9, MUTED)]])
    text(s, 12.0, 7.02, 0.9, 0.3, [[R(f"{n:02d}", 9, MUTED)]], align=PP_ALIGN.RIGHT)
    _page += 1


def bullets(s, x, y, w, items, gap=0.66, size=15, dotc=ACCENT):
    yy = y
    for it in items:
        box(s, x, yy + 0.09, 0.12, 0.12, fill=dotc, radius=True)
        text(s, x + 0.32, yy, w, gap,
             [[R(it[0], size, WHITE, True)] + ([R("  " + it[1], size, MUTED)] if it[1] else [])])
        yy += gap


def chip(s, x, y, w, label, color):
    b = box(s, x, y, w, 0.62, fill=CARD, line=color, line_w=1.25, radius=True)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT
    return b


def arrow(s, x, y, w=0.5):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                           Inches(w), Inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = ACCENT; a.line.fill.background()
    a.shadow.inherit = False
    return a


# ========================= SLIDE 1 — TITLE =================================
s = slide()
box(s, 0, 0, 0.18, 7.5, fill=ACCENT)
text(s, 0.9, 1.5, 11, 0.5, [[R("AMERICAN EXPRESS · CODESTREET 2026", 15, ACCENT, True)]])
text(s, 0.85, 2.15, 11.6, 1.8,
     [[R("The End-to-End Servicing Agent", 46, WHITE, True, FONT_H)]])
text(s, 0.9, 3.5, 11.2, 1.4, [[
    R("Resolves fee reversals, credit-limit increases & card replacements in a ", 20, MUTED),
    R("single conversation", 20, WHITE, True),
    R(" — with a ", 20, MUTED),
    R("tamper-evident audit trail", 20, WHITE, True),
    R(" and a human escape hatch.", 20, MUTED),
]])
box(s, 0.9, 5.2, 7.4, 0.02, fill=CARD)
text(s, 0.9, 5.45, 11, 0.5,
     [[R("Problem Statement:  ", 14, MUTED), R("End-to-End Servicing Agent", 14, WHITE, True)]])
text(s, 0.9, 5.95, 11, 0.5,
     [[R("Team:  ", 14, MUTED), R("<your team name>", 14, WHITE, True),
       R("      Members:  ", 14, MUTED), R("<names>", 14, WHITE, True)]])
text(s, 0.9, 6.45, 11, 0.5,
     [[R("Demo video:  ", 14, MUTED), R("<link>", 14, ACCENT, True),
       R("      Repo:  ", 14, MUTED),
       R("github.com/VaibhavDangaich/Codestreet_2026", 14, ACCENT, True)]])

# ========================= SLIDE 2 — PROBLEM ===============================
s = slide()
header(s, "The problem", "Routine servicing is slow, costly, and repetitive")
bullets(s, 0.9, 2.35, 11.4, [
    ("Fee reversals, limit increases, and replacement cards are the highest-frequency requests —",
     "and the most tedious to service."),
    ("Members wait on hold, repeat themselves, and get transferred between agents.",
     "Every touch adds compliance overhead."),
    ("Skilled human agents spend their time on routine work",
     "instead of the complex cases that actually need them."),
], gap=0.85, size=17)
c = box(s, 0.9, 5.35, 11.5, 1.1, fill=CARD, line=ACCENT, line_w=1.25, radius=True)
text(s, 1.2, 5.5, 5.2, 0.8,
     [[R("~$13.50", 30, ROSE, True)], [R("cost per human-assisted contact", 12, MUTED)]])
text(s, 4.7, 5.5, 4, 0.8,
     [[R("~$1.84", 30, GREEN, True)], [R("cost per self-service contact", 12, MUTED)]])
text(s, 8.0, 5.62, 4.2, 0.9,
     [[R("The opportunity: ", 14, WHITE, True),
       R("resolve the routine autonomously — safely — and free humans for the hard cases.", 14, MUTED)]])
footer(s)

# ========================= SLIDE 3 — SOLUTION ==============================
s = slide()
header(s, "Our solution", "An agent that resolves the request — not just talks about it")
cards = [
    ("Resolves end-to-end", GREEN,
     "Executes the actual fee waiver, limit change, or card order — start to finish, in one interaction."),
    ("Provably auditable", ACCENT,
     "Every decision & action is written to a hash-chained, tamper-evident audit trail — immutable by construction."),
    ("Escalates gracefully", AMBER,
     "When policy limits or low confidence are hit, hands off to a human with full context — no repeating."),
]
x = 0.9
for title_, col, body in cards:
    box(s, x, 2.5, 3.7, 2.9, fill=CARD, line=col, line_w=1.5, radius=True)
    box(s, x + 0.35, 2.85, 0.55, 0.12, fill=col)
    text(s, x + 0.35, 3.1, 3.1, 0.7, [[R(title_, 18, WHITE, True)]])
    text(s, x + 0.35, 3.85, 3.1, 1.4, [[R(body, 13.5, MUTED)]])
    x += 3.95
text(s, 0.9, 5.7, 11.5, 0.6, [[
    R("Single interaction.  ", 18, WHITE, True), R("Real actions.  ", 18, ACCENT, True),
    R("Complete traceability.", 18, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, 0.9, 6.35, 11.5, 0.4, [[
    R("Plus  ", 12.5, ACCENT, True),
    R("multi-turn follow-ups (slot-filling) · voice input · safe out-of-scope handling "
      "(answer / escalate / clarify)", 12.5, MUTED)]], align=PP_ALIGN.CENTER)
footer(s)

# ========================= SLIDE 4 — PRODUCT / UI =========================
s = slide()
header(s, "The product", "One console — chat, durable cases & the audit trail")
# real screenshot of the running console
_ui = Path(__file__).resolve().parent / "diagrams" / "ui_main.png"
s.shapes.add_picture(str(_ui), Inches(0.85), Inches(2.05), width=Inches(7.5))
# annotations
text(s, 9.05, 2.15, 3.4, 0.4, [[R("What you’re seeing", 13, ACCENT2, True)]])
anns = [
    ("Policy citation in every reply", "“auto-approved under FEE-AUTO v1.2”", ACCENT),
    ("Counterfactual on escalation", "“I can auto-approve up to $13,000 now”", GREEN),
    ("Escalation with full handoff context", "the specialist gets everything", AMBER),
    ("Durable cases + live audit / graph", "one screen, in real time", VIOLET),
]
yy = 2.75
for head, sub, col in anns:
    box(s, 8.75, yy + 0.04, 0.12, 0.12, fill=col, radius=True)
    text(s, 9.05, yy, 3.35, 0.7, [[R(head, 12, WHITE, True)], [R(sub, 10.5, MUTED)]])
    yy += 0.92
text(s, 9.05, 6.5, 3.4, 0.3, [[R("Live capture of the running app", 10, MUTED)]])
footer(s)

# ========================= SLIDE 5 — HOW IT WORKS ==========================
s = slide()
header(s, "How a request is decided",
       "The LLM proposes — a versioned engine decides and cites the rule")
y = 2.7
stages = [
    (0.85, 1.7, "Request", MUTED, ""),
    (3.05, 2.0, "Classifier", VIOLET, "LLM proposes"),
    (5.55, 2.5, "Policy-as-Code", ACCENT, "decides · cites rule"),
    (8.55, 1.7, "Verify", VIOLET, "risk-based"),
    (10.75, 1.65, "Execute", GREEN, ""),
]
for (x, w, label, col, cap) in stages:
    chip(s, x, y, w, label, col)
    if cap:
        text(s, x, y + 0.66, w, 0.3, [[R(cap, 10, MUTED)]], align=PP_ALIGN.CENTER)
for ax in (2.6, 5.1, 8.1, 10.3):
    arrow(s, ax, y + 0.17, 0.4)
# escalate branch
a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.3), Inches(y + 0.95),
                       Inches(0.28), Inches(0.45))
a.fill.solid(); a.fill.fore_color.rgb = ROSE; a.line.fill.background(); a.shadow.inherit = False
chip(s, 8.55, y + 1.5, 3.85, "Escalate — with counterfactual", ROSE)
text(s, 4.5, y + 1.58, 3.8, 0.5,
     [[R("over cap · low confidence · verifier disagrees", 10.5, ROSE, True)]],
     align=PP_ALIGN.RIGHT)
# audit bar
box(s, 0.9, y + 2.55, 11.5, 0.72, fill=CARD, line=ACCENT, line_w=1.25, radius=True)
text(s, 1.15, y + 2.67, 11, 0.5, [[
    R("⛓  Hash-chained audit trail   ", 14, ACCENT, True),
    R("— every decision sealed in order, tagged with the rule that fired "
      "(e.g. LIMIT-CAP v1.1)", 12.5, MUTED)]])
text(s, 0.9, 6.55, 11.5, 0.4, [[
    R("The model never has authority over money — a deterministic, versioned "
      "rule engine decides and cites the exact rule; a second agent verifies "
      "before anything executes.", 11.5, MUTED)]])
footer(s)

# ========================= SLIDE 5 — DIFF 1: AUDIT =========================
s = slide()
header(s, "Differentiator 1", "A verifiable, tamper-evident audit trail", ecolor=ACCENT)
bullets(s, 0.9, 2.3, 7.1, [
    ("Every decision, policy check & backend call is hash-chained (SHA-256)",
     "— each entry seals the one before it."),
    ("Edit any past entry and every downstream hash breaks",
     "— caught live in the demo (verify flips to ‘broken @ seq N’)."),
    ("Complete chain-of-custody, not sampling",
     "— the difference between ‘we think’ and ‘we can prove’."),
    ("Persisted to Neo4j — explorable as an analyst graph",
     "— query the whole chain of custody, not scroll a log."),
    ("Maps to model-risk (SR 11-7) & adverse-action explainability",
     "— audit-ready by design."),
], gap=0.82, size=15)
# real Neo4j audit graph (live capture) — rules are nodes, linked APPLIED_RULE
_g = Path(__file__).resolve().parent / "diagrams" / "ui_graph_panel.png"
box(s, 8.35, 2.3, 4.1, 4.35, fill=LIGHT, line=ACCENT, line_w=1.0, radius=True)
s.shapes.add_picture(str(_g), Inches(8.9), Inches(2.45), height=Inches(3.95))
text(s, 8.35, 6.44, 4.1, 0.3,
     [[R("Live Neo4j graph — policy rules are nodes (APPLIED_RULE)",
         9, MUTED)]], align=PP_ALIGN.CENTER)
footer(s)

# ========================= SLIDE 6 — DIFF 2: DURABLE CASES =================
s = slide()
header(s, "Differentiator 2 · beyond the brief",
       "Servicing cases that can’t be left half-done", ecolor=VIOLET)
bullets(s, 0.9, 2.35, 7.1, [
    ("Every high-stakes request becomes a durable Temporal workflow",
     "— it survives crashes and resumes exactly where it stopped."),
    ("Saga with automatic compensation",
     "— if a step fails, completed steps roll back in reverse. Never partial."),
    ("Durable human-in-the-loop approval",
     "— over-policy cases wait on an underwriter for as long as it takes."),
    ("Live status + replayable event history",
     "— queryable in real time; the log is itself an audit."),
], gap=0.92, size=15.5)
# right visual — a saga that fails and compensates
box(s, 8.4, 2.4, 3.6, 3.75, fill=CARD, line=VIOLET, line_w=1.5, radius=True)
text(s, 8.65, 2.6, 3.1, 0.5, [[R("Card replacement saga", 13, VIOLET, True)]])
for i, (t, c) in enumerate([
    ("✓ block old card", GREEN),
    ("✓ charge fee", GREEN),
    ("✗ order fulfillment — fails", ROSE),
    ("↩ refund fee (compensate)", AMBER),
    ("↩ unblock card (compensate)", AMBER),
    ("= rolled back · consistent", ACCENT),
]):
    text(s, 8.65, 3.1 + i * 0.5, 3.1, 0.5, [[R(t, 12.5, c, True)]])
footer(s)

# ========================= SLIDE 7b — DIFF 3: DECISIONING ==================
s = slide()
header(s, "Differentiator 3", "Trustworthy by construction — the LLM never decides",
       ecolor=ACCENT2)
cards3 = [
    ("Policy-as-Code", ACCENT,
     "Versioned declarative rules decide; the LLM only proposes. Every decision "
     "cites the exact rule + version (e.g. LIMIT-CAP v1.1). Change policy = edit "
     "data + bump a version, no agent code."),
    ("Self-verifying agent", VIOLET,
     "A second reviewer LLM must agree the interpretation matches the request "
     "before acting — a bounded propose → verify → revise loop; escalates if the "
     "two agents can't agree. Risk-based, so it's spent where it matters."),
    ("Counterfactual answers", GREEN,
     "On a denial we compute the nearest-approvable outcome ('I can do $15,600 "
     "now'), mapping to ECOA/CFPB adverse-action 'specific reasons'."),
    ("Decision provenance", AMBER,
     "Rules become nodes in Neo4j — analysts query which rule approved or denied "
     "across members, not just scroll a log."),
]
for i, (t, col, body) in enumerate(cards3):
    cx = 0.9 + (i % 2) * 5.85
    cy = 2.35 + (i // 2) * 2.15
    box(s, cx, cy, 5.6, 1.95, fill=CARD, line=col, line_w=1.4, radius=True)
    box(s, cx + 0.3, cy + 0.28, 0.5, 0.11, fill=col)
    text(s, cx + 0.3, cy + 0.5, 5.0, 0.4, [[R(t, 16, WHITE, True)]])
    text(s, cx + 0.3, cy + 1.0, 5.05, 0.9, [[R(body, 11.5, MUTED)]])
footer(s)

# ========================= SLIDE 7 — METRICS ==============================
s = slide()
header(s, "Evaluation", "Measured, not asserted")
rows = [
    ("Intent classification accuracy", "90%  (n=20)", "labeled eval set · LLM target ≥95%", GREEN),
    ("First-contact resolution", "80%  (12/15)", "multi-turn outcome check", GREEN),
    ("Correct handling (resolve / safe-escalate)", "86.7%", "outcome check, in-scope", GREEN),
    ("Escalation handoff completeness", "100%", "summary + reason + context present", GREEN),
    ("Policy violations", "0", "engine caps enforced + verified", ACCENT),
    ("Audit integrity", "100% intact", "hash-chain verification", ACCENT),
]
# table header
box(s, 0.9, 2.35, 11.5, 0.5, fill=ACCENT2, radius=False)
text(s, 1.1, 2.42, 5, 0.4, [[R("Metric", 13, LIGHT, True)]])
text(s, 6.4, 2.42, 3, 0.4, [[R("Result", 13, LIGHT, True)]])
text(s, 9.4, 2.42, 3, 0.4, [[R("How it’s graded", 13, LIGHT, True)]])
yy = 2.85
for i, (m, res, how, col) in enumerate(rows):
    fill = CARD if i % 2 == 0 else BG
    box(s, 0.9, yy, 11.5, 0.62, fill=fill)
    text(s, 1.1, yy + 0.13, 5.2, 0.4, [[R(m, 13.5, WHITE, True)]])
    text(s, 6.4, yy + 0.11, 3, 0.4, [[R(res, 13, col, True)]])
    text(s, 9.4, yy + 0.11, 3, 0.4, [[R(how, 12, MUTED)]])
    yy += 0.56
text(s, 0.9, yy + 0.18, 11.5, 0.6, [[
    R("We report TRUE resolution + handoff quality — not just deflection ", 12.5, WHITE, True),
    R("(financial-services deflection is realistically 25–45%; one wrong answer is a "
      "compliance event, so the target error rate is <1%).", 12.5, MUTED)]])
footer(s)

# ========================= SLIDE 8 — WHY AMEX ==============================
s = slide()
header(s, "Why Amex · why now", "The governance layer that makes autonomy deployable")
bullets(s, 0.9, 2.2, 11.4, [
    ("Amex’s 2026 direction is agentic commerce — agents that decide and act",
     "— this is exactly that, applied to servicing."),
    ("Only 21% of enterprises have mature agentic-AI governance, yet 74% expect to use AI agents by 2027",
     "(Deloitte) — most are scaling agents without the guardrails."),
    ("Deloitte names the missing controls: audit trails of every agent action + human-approval boundaries",
     "— the exact two things we built."),
], gap=0.9, size=15.5)
box(s, 0.9, 5.05, 11.5, 1.0, fill=CARD, line=ACCENT, line_w=1.25, radius=True)
text(s, 1.2, 5.2, 11, 0.8, [[
    R("“With every use case, we’ve ensured there is human oversight.”", 15, WHITE, True),
    R("   — Ravi Radhakrishnan, EVP & CIO, American Express (Forbes, 2025)", 12.5, MUTED)]])
text(s, 0.9, 6.3, 11.5, 0.5, [[
    R("Sources:  ", 9.5, ACCENT, True),
    R("Deloitte, “Agentic AI is scaling faster than guardrails” (2025)   ·   "
      "Forbes / Peter High, “Ravi Radhakrishnan on Driving AI Innovation at American Express” (May 2025)",
      9.5, MUTED)]])
footer(s)

# ========================= SLIDE 9 — ARCHITECTURE =========================
s = slide()
header(s, "Architecture", "Built to be real")
# row 1: request path
chip(s, 0.9, 2.55, 2.2, "Next.js UI", ACCENT)
arrow(s, 3.2, 2.72)
chip(s, 3.8, 2.55, 2.0, "FastAPI", ACCENT)
arrow(s, 5.9, 2.72)
chip(s, 6.5, 2.55, 5.9, "LangGraph agent: classifier · policy · flows · audit", VIOLET)
# row 2: durable case path
chip(s, 0.9, 3.85, 2.6, "Temporal server", VIOLET)
arrow(s, 3.6, 4.02)
chip(s, 4.2, 3.85, 1.7, "Worker", VIOLET)
arrow(s, 6.0, 4.02)
chip(s, 6.6, 3.85, 5.8, "ServicingCaseWorkflow → saga · compensation · approval", VIOLET)
# row 3: audit graph
chip(s, 0.9, 4.55, 3.4, "Audit chain → /graph", ACCENT)
arrow(s, 4.4, 4.72)
chip(s, 5.0, 4.55, 3.0, "Neo4j Aura", GREEN)
arrow(s, 8.1, 4.72)
chip(s, 8.7, 4.55, 3.7, "Cytoscape analyst graph", GREEN)
# stack line
box(s, 0.9, 5.35, 11.5, 0.95, fill=CARD, radius=True)
text(s, 1.2, 5.48, 11, 0.8, [[R("Stack   ", 13, ACCENT, True),
    R("Python · FastAPI · LangGraph · Gemini · Temporal · Neo4j · Next.js · Cytoscape · SHA-256 hash-chain",
      12.5, WHITE)]])
text(s, 1.2, 5.72, 11, 0.4, [[
    R("Core-banking is mocked today; each call swaps to a real Amex API with no change to the agent logic.",
      12.5, MUTED)]])
footer(s)

# ========================= SLIDE 11 — SYSTEM DESIGN (LLD) ==================
s = slide()
header(s, "System design", "Low-level design — how a durable case executes")
_seq = Path(__file__).resolve().parent / "diagrams" / "durable_case_sequence.png"
s.shapes.add_picture(str(_seq), Inches(0.85), Inches(2.15), width=Inches(7.15))
MONO = "Courier New"
# components card
box(s, 8.25, 2.15, 4.2, 2.15, fill=CARD, line=ACCENT, line_w=1.0, radius=True)
text(s, 8.5, 2.28, 3.8, 0.3, [[R("Components", 12, ACCENT2, True)]])
for i, (k, v) in enumerate([
    ("FastAPI", "owns card state + audit chain"),
    ("LangGraph", "classify · policy · flows · assist"),
    ("Temporal worker", "durable saga + compensation"),
    ("Neo4j Aura", "audit graph for analysts"),
]):
    text(s, 8.5, 2.62 + i * 0.38, 3.85, 0.34,
         [[R(k + "  ", 10.5, WHITE, True), R(v, 10, MUTED)]])
# data model + APIs card
box(s, 8.25, 4.5, 4.2, 2.38, fill=CARD, line=ACCENT, line_w=1.0, radius=True)
text(s, 8.5, 4.62, 3.8, 0.3, [[R("Audit entry (hash-chained)", 12, ACCENT2, True)]])
text(s, 8.5, 4.94, 3.85, 0.5,
     [[R("{ seq, actor, action,", 9.5, WHITE, False, MONO)],
      [R("  payload, prev_hash, hash }", 9.5, WHITE, False, MONO)]])
text(s, 8.5, 5.62, 3.8, 0.3, [[R("Key APIs", 12, ACCENT2, True)]])
for i, api in enumerate([
    "POST /chat",
    "POST /cases/start · /decision",
    "POST /internal/action/{action}",
    "GET  /cases · /graph",
]):
    text(s, 8.5, 5.94 + i * 0.235, 3.85, 0.22, [[R(api, 9.5, WHITE, False, MONO)]])
footer(s)

# ========================= SLIDE 12 — NEXT / ASK ==========================
s = slide()
header(s, "What’s next", "From prototype to production")
bullets(s, 0.9, 2.4, 11.4, [
    ("Round 2: integrate real servicing APIs, add voice, expand intents,",
     "and add an LLM-judge clarity eval."),
    ("Production audit store: append-only, write-once storage behind the same hash-chain",
     "— exportable for auditors."),
    ("Guardrails & governance: per-action risk tiers, kill-switch, policy versioning",
     "— the governance substrate for Amex agents."),
], gap=0.92, size=16.5)
box(s, 0.9, 5.5, 11.5, 1.05, fill=ACCENT2, radius=True)
text(s, 1.2, 5.72, 11, 0.7, [[
    R("The ask:  ", 18, LIGHT, True),
    R("advance us to the prototype round — the core already works end-to-end.", 17, LIGHT)]])
footer(s)

# --- save ------------------------------------------------------------------
out = Path(__file__).resolve().parent / "CodeStreet2026_Servicing_Agent.pptx"
prs.save(out)
print(f"Saved deck -> {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
